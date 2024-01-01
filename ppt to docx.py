import os
import tkinter as tk
from tkinter import filedialog
from pptx import Presentation
from docx import Document
from docx.shared import Pt
from docx.oxml.ns import qn
from docx.shared import Pt
from docx.oxml.ns import qn
from docx.enum.style import WD_STYLE_TYPE
import tkinter.messagebox as mb



def clean_text(text):
    """
    清除文本中的不兼容XML字符，包括NULL字节和控制字符。
    """
    if not text:
        return text
    # 移除NULL字节和其他非打印字符
    cleaned_text = ''.join(char for char in text if char.isprintable())
    # 确保文本为字符串类型
    cleaned_text = str(cleaned_text)
    return cleaned_text

def extract_text_from_pptx(pptx_file):
    prs = Presentation(pptx_file)
    text_runs = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                # 清理并添加文本
                cleaned_text = clean_text(paragraph.text)
                text_runs.append(cleaned_text)

    return '\n'.join(text_runs)

def set_default_font(document, font_name='宋体', font_size=12):
    # 设置文档默认字体
    style = document.styles['Normal']
    font = style.font
    font.name = font_name
    font.size = Pt(font_size)
    font._element.rPr.rFonts.set(qn('w:eastAsia'), font_name)

def add_paragraph_with_font(document, text):
    paragraph = document.add_paragraph(text)
    return paragraph

def process_pptx_files_in_folder(folder_path, output_file):
    document = Document()
    set_default_font(document)  # 设置文档默认字体为宋体

    for filename in os.listdir(folder_path):
        if filename.lower().endswith(('.pptx', '.ppt')):
            file_path = os.path.join(folder_path, filename)
            extracted_text = extract_text_from_pptx(file_path)
            add_paragraph_with_font(document, extracted_text)
            document.add_page_break()

    document.save(output_file)


def select_folder_and_process():
    folder_path = filedialog.askdirectory()
    if folder_path:
        output_file = filedialog.asksaveasfilename(
            defaultextension=".docx",
            filetypes=[("Word Documents", "*.docx")],
            title="Save the document as"
        )
        if output_file:
            process_pptx_files_in_folder(folder_path, output_file)
            mb.showinfo("转换完成", "PPT到Word的转换已完成！")

root = tk.Tk()
root.title("PPT to Word Converter")
root.geometry("400x200")

btn_select_folder = tk.Button(root, text="Select Folder and Convert", command=select_folder_and_process)
btn_select_folder.pack(pady=50)

root.mainloop()
